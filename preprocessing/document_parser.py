# preprocessing/document_parser.py

import os
import logging
import subprocess
from typing import Dict, List, Union, Tuple
import fitz  # PyMuPDF
import docx
import olefile
import pptx
import pandas as pd
import io

try:
    from openpyxl_image_loader import SheetImageLoader
    OPENPYXL_IMAGE_LOADER_INSTALLED = True
except ImportError:
    OPENPYXL_IMAGE_LOADER_INSTALLED = False

log = logging.getLogger(__name__)

def parse_document(file_path: str) -> Dict[str, Union[str, List[Tuple[int, bytes]]]]:
    ext = os.path.splitext(file_path)[1].lower()
    
    parser_map = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".hwp": _parse_hwp,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
    }
    
    parser_func = parser_map.get(ext)
    
    if parser_func:
        return parser_func(file_path)
    else:
        log.warning(f"지원하지 않는 문서 파일 형식입니다: {file_path}")
        return {"text": "", "images": []}

def _parse_pdf(file_path: str) -> Dict[str, Union[str, List[Tuple[int, bytes]]]]:
    try:
        doc = fitz.open(file_path)
        full_text = ""
        images = []
        for page in doc:
            full_text += page.get_text() + "\n"
            for img in page.get_images(full=True):
                xref = img[0]
                base_image = doc.extract_image(xref)
                images.append((len(images) + 1, base_image["image"]))
        doc.close()
        return {"text": full_text, "images": images}
    except Exception as e:
        log.error(f"PDF 파일 처리 중 오류 발생: {os.path.basename(file_path)}", exc_info=True)
        return {"text": "", "images": []}

def _parse_docx(file_path: str) -> Dict[str, Union[str, List[Tuple[int, bytes]]]]:
    try:
        doc = docx.Document(file_path)
        full_text = "\n".join([para.text for para in doc.paragraphs])
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                images.append((len(images) + 1, rel.target_part.blob))
        return {"text": full_text, "images": images}
    except Exception as e:
        log.error(f"DOCX 파일 처리 중 오류 발생: {os.path.basename(file_path)}", exc_info=True)
        return {"text": "", "images": []}
    
def _parse_hwp(file_path: str) -> Dict[str, Union[str, List[Tuple[int, bytes]]]]:
    full_text = ""
    images = []
    file_name = os.path.basename(file_path)

    try:
        # --- 1. 텍스트 추출 (subprocess 사용) ---
        log.info(f"hwp5txt 명령어를 사용하여 텍스트 추출 시도: {file_name}")
        # hwp5txt <파일경로> 명령을 실행하고, 결과를 stdout으로 받음
        result = subprocess.run(
            ['hwp5txt', file_path],
            capture_output=True,  # 표준 출력/오류를 캡처
            text=True,            # 출력을 텍스트(문자열)로 디코딩
            check=True,           # 반환 코드가 0이 아니면 예외 발생
            encoding='utf-8'      # 인코딩 명시
        )
        full_text = result.stdout
        log.info(f"hwp5txt 텍스트 추출 성공.")

    except FileNotFoundError:
        log.error("'hwp5txt' 명령어를 찾을 수 없습니다. 'pip install pyhwp'가 올바르게 설치되었는지 확인해주세요.")
        # 텍스트 추출에 실패해도 이미지 추출은 시도하도록 넘어감
    except subprocess.CalledProcessError as e:
        log.error(f"hwp5txt 명령어 실행 중 오류 발생: {file_name}\nError: {e.stderr}")
        # 텍스트 추출에 실패해도 이미지 추출은 시도하도록 넘어감
    except Exception as e:
        log.error(f"HWP 텍스트 추출 중 예상치 못한 오류 발생: {file_name}", exc_info=True)

    try:
        # --- 2. 이미지 추출 (olefile 사용) ---
        log.info(f"olefile을 사용하여 이미지 추출 시도: {file_name}")
        with olefile.OleFileIO(file_path) as ole:
            for stream_path in ole.listdir():
                # HWP 이미지는 'Pictures' 스트림 디렉토리 아래에 저장됨
                if stream_path and stream_path[0] == 'Pictures':
                    image_data = ole.openstream(stream_path).read()
                    images.append((len(images) + 1, image_data))
        if images:
            log.info(f"{len(images)}개의 이미지 추출 성공.")
    
    except Exception as e:
        # olefile 라이브러리가 파일을 열 수 없는 경우 (예: 파일 손상)
        log.error(f"HWP에서 olefile로 이미지 추출 중 오류 발생: {file_name}", exc_info=True)

    return {"text": full_text, "images": images}

def _parse_pptx(file_path: str) -> Dict[str, Union[str, List[Tuple[int, bytes]]]]:
    try:
        presentation = pptx.Presentation(file_path)
        full_text_list, images = [], []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text_list.append(shape.text_frame.text)
                if isinstance(shape, pptx.shapes.picture.Picture):
                    images.append((len(images) + 1, shape.image.blob))
        return {"text": "\n\n".join(full_text_list), "images": images}
    except Exception as e:
        log.error(f"PPTX 파일 처리 중 오류 발생: {os.path.basename(file_path)}", exc_info=True)
        return {"text": "", "images": []}

def _parse_xlsx(file_path: str) -> Dict[str, Union[str, List[Tuple[int, bytes]]]]:
    try:
        full_text_list, images = [], []
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, header=None)
            if not df.empty:
                full_text_list.append(f"--- 시트: {sheet_name} ---\n{df.to_string(index=False)}")
        
        if OPENPYXL_IMAGE_LOADER_INSTALLED:
            image_loader = SheetImageLoader(file_path)
            for sheet_name in image_loader.sheet_names:
                for image_name in image_loader.get_images_in_sheet(sheet_name):
                    image = image_loader.get_image(image_name)
                    with io.BytesIO() as img_byte_arr:
                        image.save(img_byte_arr, format='PNG')
                        images.append((len(images) + 1, img_byte_arr.getvalue()))
        else:
            log.warning("엑셀 이미지 추출을 위해 'pip install openpyxl-image-loader'를 설치해주세요. 이미지 추출을 건너뜁니다.")

        return {"text": "\n\n".join(full_text_list), "images": images}
    except Exception as e:
        log.error(f"XLSX 파일 처리 중 오류 발생: {os.path.basename(file_path)}", exc_info=True)
        return {"text": "", "images": []}